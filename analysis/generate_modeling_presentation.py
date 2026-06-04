from __future__ import annotations

from pathlib import Path
import sys

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from generate_kickoff_presentation import (
    BLUE,
    GREEN,
    INK,
    MUTED,
    ORANGE,
    PAPER,
    RED,
    TEAL,
    WHITE,
    add_background,
    add_bullets,
    add_card,
    add_footer,
    add_image,
    add_textbox,
    add_title,
    set_fill,
    set_run,
)


PACKAGE_DIR = ROOT / "sgRNA_synthesizability"
SLIDES_MD = PACKAGE_DIR / ".github" / "slides.md"
OUTPUT_DIR = ROOT / "analysis" / "modeling_outputs"
FEATURE_PLOT_DIR = ROOT / "analysis" / "feature_target_plots"
FEATURE_DIR = ROOT / "analysis" / "features"
OUTPUT = ROOT / "sgRNA_synthesizability_modeling_summary.pptx"
MAINPEAK_FIT_PLOT = OUTPUT_DIR / "presentation_mainpeak_observed_predicted_fit.png"

FEATURE_LABELS = {
    "rna_count_middle": "RNA Count in Middle Region",
    "modified_token_position_std": "Standard Deviation of Modified Token Positions",
    "token_rC_count": "rC Token Count",
    "lna_span_fraction": "LNA Span Fraction",
    "base_U_count": "U Base Count",
    "star_count_3p": "3' Star Modification Count",
    "token_lG_fraction": "lG Token Fraction",
    "longest_rU_run": "Longest rU Run",
    "longest_lC_run": "Longest lC Run",
}


def require_inputs() -> None:
    required = [
        SLIDES_MD,
        OUTPUT_DIR / "presentation_summary.md",
        OUTPUT_DIR / "learning_curve.png",
        OUTPUT_DIR / "mainpeak_regression_oof_predictions.csv",
        OUTPUT_DIR / "mainpeak_regression_derived_classifier_metrics.csv",
        OUTPUT_DIR / "mainpeak_regression_derived_classifier_class_metrics.csv",
    ]
    missing = [path for path in required if not path.exists()]
    if missing:
        missing_text = "\n".join(str(path.relative_to(ROOT)) for path in missing)
        raise FileNotFoundError(f"Missing presentation input assets:\n{missing_text}")


def add_metric(slide, left, top, width, height, label, value, note=None, fill=RGBColor(238, 244, 235)):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(card, fill, RGBColor(218, 211, 198))
    value_box = add_textbox(slide, left + Inches(0.14), top + Inches(0.15), width - Inches(0.28), Inches(0.38))
    paragraph = value_box.text_frame.paragraphs[0]
    paragraph.text = value
    paragraph.alignment = PP_ALIGN.CENTER
    set_run(paragraph.runs[0], size=20, color=INK, bold=True)
    label_box = add_textbox(slide, left + Inches(0.14), top + Inches(0.57), width - Inches(0.28), Inches(0.38))
    paragraph = label_box.text_frame.paragraphs[0]
    paragraph.text = label
    paragraph.alignment = PP_ALIGN.CENTER
    set_run(paragraph.runs[0], size=10, color=MUTED, bold=True)
    if note:
        note_box = add_textbox(slide, left + Inches(0.18), top + Inches(0.95), width - Inches(0.36), height - Inches(1.02))
        paragraph = note_box.text_frame.paragraphs[0]
        paragraph.text = note
        paragraph.alignment = PP_ALIGN.CENTER
        set_run(paragraph.runs[0], size=9, color=MUTED)


def add_result_row(slide, y, endpoint, model, performance, color):
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), y + Inches(0.08), Inches(0.12), Inches(0.52))
    set_fill(marker, color)
    endpoint_box = add_textbox(slide, Inches(0.95), y, Inches(2.45), Inches(0.72))
    paragraph = endpoint_box.text_frame.paragraphs[0]
    paragraph.text = endpoint
    set_run(paragraph.runs[0], size=13, color=INK, bold=True)
    model_box = add_textbox(slide, Inches(3.55), y, Inches(3.1), Inches(0.72))
    paragraph = model_box.text_frame.paragraphs[0]
    paragraph.text = model
    set_run(paragraph.runs[0], size=12, color=INK)
    perf_box = add_textbox(slide, Inches(6.95), y, Inches(5.65), Inches(0.72))
    paragraph = perf_box.text_frame.paragraphs[0]
    paragraph.text = performance
    set_run(paragraph.runs[0], size=12, color=INK)


def add_performance_table(slide, left, top, width, height):
    columns = ["Endpoint", "Best current model / feature set", "PR-AUC", "ROC AUC", "F1 score"]
    rows = [
        ("Source Pass/Fail", "71-feature correlation-pruned HGB", "0.1139", "0.6529", "0.0800"),
        ("Full-feature Pass/Fail", "161-feature HGB", "0.1087", "0.6304", "0.1165"),
        ("Formula Need_Review", "71-feature correlation-pruned HGB", "0.4082", "0.7719", "0.3467"),
        ("MainPeak-derived Pass/Fail", "Elastic-net MainPeak regressor + rule threshold", "0.1069", "0.6217", "0.1212"),
        ("Compact 9-feature Pass/Fail", "9-feature L1 logistic panel", "0.1078", "0.5483", "0.0693"),
        ("Compact 9-feature Need_Review", "9-feature HGB panel", "0.3749", "0.7677", "0.3529"),
    ]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(2.15)
    table.columns[1].width = Inches(4.05)
    table.columns[2].width = Inches(1.55)
    table.columns[3].width = Inches(1.55)
    table.columns[4].width = Inches(1.75)
    for col_index, column in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = column
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 48, 60)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        set_run(paragraph.runs[0], size=10, color=WHITE, bold=True)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else RGBColor(244, 240, 232)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if col_index < 2 else PP_ALIGN.CENTER
            set_run(paragraph.runs[0], size=9, color=INK, bold=(col_index == 0))
    return table_shape


def add_feature_chip(slide, text, left, top, width, color):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.34))
    set_fill(chip, color)
    frame = chip.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    set_run(paragraph.runs[0], size=9, color=WHITE, bold=True)


def save_mainpeak_fit_plot() -> Path:
    predictions = pd.read_csv(OUTPUT_DIR / "mainpeak_regression_oof_predictions.csv")
    plot_data = predictions[predictions["model"] == "extra_trees"].copy()
    plot_data["observed_mainpeak"] = pd.to_numeric(plot_data["observed_mainpeak"], errors="coerce")
    plot_data["predicted_mainpeak"] = pd.to_numeric(plot_data["predicted_mainpeak"], errors="coerce")
    plot_data = plot_data.dropna(subset=["observed_mainpeak", "predicted_mainpeak"])

    slope, intercept = np.polyfit(plot_data["observed_mainpeak"], plot_data["predicted_mainpeak"], 1)
    x_values = np.linspace(plot_data["observed_mainpeak"].min(), plot_data["observed_mainpeak"].max(), 100)
    y_values = slope * x_values + intercept
    correlation = plot_data["observed_mainpeak"].corr(plot_data["predicted_mainpeak"])

    plt.rcParams.update(
        {
            "figure.facecolor": "#faf8f2",
            "axes.facecolor": "#faf8f2",
            "axes.edgecolor": "#5b697a",
            "axes.labelcolor": "#182330",
            "xtick.color": "#182330",
            "ytick.color": "#182330",
            "text.color": "#182330",
            "font.family": "DejaVu Sans",
        }
    )
    fig, ax = plt.subplots(figsize=(7.2, 4.6))
    ax.scatter(plot_data["observed_mainpeak"], plot_data["predicted_mainpeak"], s=18, alpha=0.32, color="#2670ae", edgecolors="none")
    ax.plot(x_values, y_values, color="#be3e3e", linewidth=2.4, label=f"Fitted regression line (r={correlation:.2f})")
    ax.plot(x_values, x_values, color="#5b697a", linewidth=1.4, linestyle="--", label="Perfect prediction")
    ax.set_title("Predicted MainPeak vs Observed MainPeak", fontsize=15, weight="bold", pad=12)
    ax.set_xlabel("Observed UHPLC % MainPeak")
    ax.set_ylabel("Predicted UHPLC % MainPeak")
    ax.legend(frameon=False, loc="lower left")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.22)
    fig.tight_layout()
    fig.savefig(MAINPEAK_FIT_PLOT, dpi=180)
    plt.close(fig)
    return MAINPEAK_FIT_PLOT


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(7.5))
    set_fill(band, RGBColor(22, 40, 50))
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.82), prs.slide_width, Inches(0.18))
    set_fill(ribbon, GREEN)
    title = add_textbox(slide, Inches(0.74), Inches(0.88), Inches(11.9), Inches(1.2))
    paragraph = title.text_frame.paragraphs[0]
    paragraph.text = "sgRNA Synthesizability Prediction"
    set_run(paragraph.runs[0], size=39, color=WHITE, bold=True)
    subtitle = add_textbox(slide, Inches(0.8), Inches(2.16), Inches(10.8), Inches(0.62))
    paragraph = subtitle.text_frame.paragraphs[0]
    paragraph.text = "Modeling sequence-derived risk signals for HPLC QC review prioritization"
    set_run(paragraph.runs[0], size=20, color=RGBColor(222, 232, 226))
    add_metric(slide, Inches(1.1), Inches(3.55), Inches(3.2), Inches(1.35), "modeling deliverable", "6 slides", "one title page plus five concise content pages", RGBColor(238, 244, 235))
    add_metric(slide, Inches(5.05), Inches(3.55), Inches(3.2), Inches(1.35), "input signal", "Sequence", "decorated sgRNA chemistry parsed into features", RGBColor(240, 243, 249))
    add_metric(slide, Inches(9.0), Inches(3.55), Inches(3.2), Inches(1.35), "decision support", "Review", "rank designs for QC risk follow-up", RGBColor(244, 239, 229))
    return slide


def slide_background_goals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Background, Problem, and Goals", "Scientific context, operational need, and hackathon deliverables.")
    add_card(slide, Inches(0.82), Inches(1.55), Inches(3.55), Inches(1.28), "Scientific context", "Decorated sgRNAs combine RNA sequence, base chemistry, and backbone modifications that can affect synthesis purity.", RGBColor(234, 244, 242))
    add_card(slide, Inches(4.88), Inches(1.55), Inches(3.55), Inches(1.28), "Operational context", "HPLC QC failures create review, remake, and delivery-risk decisions after designs have already entered production.", RGBColor(244, 239, 229))
    add_card(slide, Inches(8.94), Inches(1.55), Inches(3.55), Inches(1.28), "Prediction task", "Given a decorated sgRNA sequence, estimate risk of HPLC QC failure or low purity before QC is complete.", RGBColor(240, 243, 249))
    add_bullets(
        slide,
        Inches(1.0),
        Inches(3.45),
        Inches(11.1),
        Inches(1.65),
        [
            "Hackathon goal: build a reproducible feature and modeling workflow that turns decorated sequence text into presentation-ready risk evidence.",
            "Decision goal: support review prioritization and feature insight rather than automated release decisions.",
            "Scientific goal: identify interpretable sequence-chemistry patterns associated with QC and purity outcomes.",
        ],
        size=16,
    )
    return slide


def slide_data_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Data, Feature Families, and Modeling Approach", "Problem scale, feature families, and the best-performing model families.")
    add_metric(slide, Inches(0.76), Inches(1.55), Inches(2.25), Inches(1.25), "binary-labeled rows", "2,646", "source Pass/Fail endpoint", RGBColor(238, 244, 235))
    add_metric(slide, Inches(3.25), Inches(1.55), Inches(2.25), Inches(1.25), "source Fail rows", "86", "3.25% baseline failure rate", RGBColor(247, 235, 232))
    add_metric(slide, Inches(5.74), Inches(1.55), Inches(2.25), Inches(1.25), "features", "161", "sequence-derived numeric inputs", RGBColor(240, 243, 249))
    add_metric(slide, Inches(8.23), Inches(1.55), Inches(2.25), Inches(1.25), "compact panel", "9", "readable cross-endpoint signals", RGBColor(244, 239, 229))
    add_metric(slide, Inches(10.72), Inches(1.55), Inches(1.8), Inches(1.25), "CV", "5-fold", "grouped by sequence", RGBColor(234, 244, 242))
    add_bullets(
        slide,
        Inches(0.85),
        Inches(3.22),
        Inches(11.2),
        Inches(2.45),
        [
            "Feature families: sequence size, base composition, motif and run structure, modification burden, positional chemistry, and parser completeness.",
            "Best-performing classifiers: correlation-pruned Histogram Gradient Boosting for source Pass/Fail and formula-derived Need_Review.",
            "Best-performing purity model: ExtraTrees regression on full sequence-derived features; top-40 reduced features are close behind.",
        ],
        size=15,
    )
    add_card(slide, Inches(1.0), Inches(5.85), Inches(11.1), Inches(0.72), "Parser note", "+A, +C, +G, +U, and +T are interpreted as LNA nucleotides before feature extraction.", RGBColor(234, 244, 242))
    return slide


def slide_modeling_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Key Performance", "Classification-compatible metrics for the best current model or feature set by endpoint.")
    add_performance_table(slide, Inches(0.78), Inches(1.55), Inches(11.75), Inches(4.65))
    add_card(slide, Inches(0.95), Inches(6.35), Inches(11.35), Inches(0.55), "Interpretation", "Need_Review is the strongest classification endpoint; source Pass/Fail remains a rare-event ranking problem.", RGBColor(244, 239, 229))
    return slide


def slide_feature_and_purity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Feature Insights and Purity Signal", "MainPeak prediction gives a quantitative view of sequence-derived purity signal.")
    add_image(slide, MAINPEAK_FIT_PLOT, Inches(0.78), Inches(1.62), Inches(5.65))
    add_bullets(
        slide,
        Inches(6.85),
        Inches(1.58),
        Inches(5.45),
        Inches(1.62),
        [
            "The observed-vs-predicted plot uses full-feature ExtraTrees MainPeak predictions with a fitted regression line and a perfect-prediction reference.",
            "MainPeak keeps quantitative purity information that binary Pass/Fail labels discard.",
            "Top-40 ExtraTrees-selected features nearly match the full model: RMSE 6.0998 vs 6.0574.",
        ],
        size=15,
    )
    chip_specs = [
        (FEATURE_LABELS["rna_count_middle"], 6.9, 3.38, 2.28, GREEN),
        (FEATURE_LABELS["modified_token_position_std"], 9.35, 3.38, 2.65, BLUE),
        (FEATURE_LABELS["token_rC_count"], 6.9, 3.88, 1.65, ORANGE),
        (FEATURE_LABELS["lna_span_fraction"], 8.75, 3.88, 1.65, TEAL),
        (FEATURE_LABELS["base_U_count"], 10.55, 3.88, 1.35, RED),
        (FEATURE_LABELS["star_count_3p"], 6.9, 4.38, 1.95, GREEN),
        (FEATURE_LABELS["token_lG_fraction"], 9.05, 4.38, 1.65, BLUE),
        (FEATURE_LABELS["longest_rU_run"], 10.88, 4.38, 1.35, ORANGE),
        (FEATURE_LABELS["longest_lC_run"], 6.9, 4.88, 1.65, TEAL),
    ]
    for text, left, top, width, color in chip_specs:
        add_feature_chip(slide, text, Inches(left), Inches(top), Inches(width), color)
    add_card(slide, Inches(6.85), Inches(5.2), Inches(5.45), Inches(1.08), "Compact interpretation", "The 9-feature panel is the best lightweight cross-endpoint explanation set, while the 71-feature HGB remains the best Pass/Fail ranking option.", RGBColor(240, 243, 249))
    return slide


def slide_recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Operational Recommendations", "Use the current models to prioritize review and guide data collection.")
    add_card(slide, Inches(0.82), Inches(1.55), Inches(3.55), Inches(1.25), "Risk ranking", "Deploy the 71-feature HGB score as a review-prioritization queue, not an automatic release decision.", RGBColor(234, 244, 242))
    add_card(slide, Inches(4.88), Inches(1.55), Inches(3.55), Inches(1.25), "Threshold setting", "Choose cutoffs from review budget or target precision, then calibrate on prospective batches.", RGBColor(244, 239, 229))
    add_card(slide, Inches(8.94), Inches(1.55), Inches(3.55), Inches(1.25), "Data growth", "Collect more true-Fail and low-MainPeak examples; these are the limiting labels.", RGBColor(247, 235, 232))
    add_image(slide, OUTPUT_DIR / "learning_curve.png", Inches(0.95), Inches(3.32), Inches(5.55))
    add_bullets(
        slide,
        Inches(7.0),
        Inches(3.35),
        Inches(5.15),
        Inches(2.5),
        [
            "Best source Pass/Fail model: 71-feature correlation-pruned HGB for top-risk enrichment.",
            "Best compact reporting view: revised 9-feature panel across Pass/Fail, Need_Review, and MainPeak.",
            "Best purity predictor: full-feature ExtraTrees, or top-40 ExtraTrees features when compactness matters.",
        ],
        size=15,
    )
    return slide


def slide_lessons_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Lessons Learned and Next Steps", "The main constraint is rare-fail evidence, not feature availability.")
    add_card(slide, Inches(0.82), Inches(1.55), Inches(3.55), Inches(1.25), "Lesson 1", "Rare source Fail labels make ranking metrics more honest than accuracy or default-threshold F1.", RGBColor(247, 235, 232))
    add_card(slide, Inches(4.88), Inches(1.55), Inches(3.55), Inches(1.25), "Lesson 2", "MainPeak and Need_Review are more learnable because they preserve or derive from stronger purity signal.", RGBColor(234, 244, 242))
    add_card(slide, Inches(8.94), Inches(1.55), Inches(3.55), Inches(1.25), "Lesson 3", "Compact panels explain cross-endpoint behavior, but 71-feature HGB remains best for source Fail ranking.", RGBColor(240, 243, 249))
    add_image(slide, OUTPUT_DIR / "learning_curve.png", Inches(0.95), Inches(3.32), Inches(5.55))
    add_bullets(
        slide,
        Inches(7.0),
        Inches(3.35),
        Inches(5.15),
        Inches(2.5),
        [
            "Calibrate thresholds against a fixed review budget or target precision.",
            "Validate on prospective or time-split batches when available.",
            "Collect more true-Fail and low-MainPeak examples to improve rare-event performance.",
        ],
        size=15,
    )
    return slide


def build_deck() -> None:
    require_inputs()
    save_mainpeak_fit_plot()
    pd.read_csv(OUTPUT_DIR / "correlation_reduced_model_comparison.csv")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        slide_title(prs),
        slide_background_goals(prs),
        slide_data_features(prs),
        slide_modeling_results(prs),
        slide_feature_and_purity(prs),
        slide_lessons_next_steps(prs),
    ]
    for number, slide in enumerate(slides, start=1):
        if number > 1:
            add_footer(slide, number)
    prs.save(OUTPUT)
    print(f"wrote {OUTPUT.relative_to(ROOT)} with {len(slides)} slides")


if __name__ == "__main__":
    build_deck()